#!/usr/bin/env python3
"""Build the results deck as a .pptx that imports cleanly into Google Slides.

Figures come from results/figures/ (written by charts.R, which reads the
runners' JSON) so no number is retyped on its way to a slide. The technique
tables are transcribed from ~/paper's Table `tab:speedup-strategies` and its
Materials and methods section, which are the record for what was done and why.

    scratchpad/pptxenv/bin/python scripts/render/deck.py   # -> results/jb2-results.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIGS = ROOT / "results" / "figures"
OUT = ROOT / "results" / "jb2-results.pptx"

INK = RGBColor(0x14, 0x1C, 0x1E)
MUTED = RGBColor(0x55, 0x6A, 0x6C)
TEAL = RGBColor(0x0E, 0x6B, 0x62)
CLAY = RGBColor(0xB0, 0x4A, 0x2C)
RULE = RGBColor(0xD5, 0xDF, 0xDE)
BAND = RGBColor(0xEE, 0xF4, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, *, bold=False, color=INK, space_after=0, first=False,
         align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return p


def rule(slide, top, left=MARGIN, width=None, color=RULE, height=Pt(1.2)):
    width = width or (W - 2 * MARGIN)
    bar = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def slide_head(slide, title, subtitle=None, eyebrow=None):
    top = MARGIN
    if eyebrow:
        tf = textbox(slide, MARGIN, top, W - 2 * MARGIN, Inches(0.26))
        para(tf, eyebrow.upper(), 11, bold=True, color=TEAL, first=True)
        top += Inches(0.32)
    tf = textbox(slide, MARGIN, top, W - 2 * MARGIN, Inches(0.62))
    para(tf, title, 30, bold=True, first=True)
    top += Inches(0.66)
    if subtitle:
        tf = textbox(slide, MARGIN, top, W - 2 * MARGIN, Inches(0.5))
        para(tf, subtitle, 14, color=MUTED, first=True)
        top += Inches(0.52)
    rule(slide, top)
    return top + Inches(0.24)


def footnote(slide, text):
    tf = textbox(slide, MARGIN, H - Inches(0.72), W - 2 * MARGIN, Inches(0.46))
    for i, line in enumerate(text.split("\n")):
        para(tf, line, 10.5, color=MUTED, first=(i == 0))


def picture(slide, name, top, *, bottom_pad=Inches(0.85)):
    """Fit a figure into the space left under the header, preserving aspect."""
    from PIL import Image  # pillow ships with the venv's python-pptx deps

    path = FIGS / name
    with Image.open(path) as im:
        iw, ih = im.size
    avail_w = W - 2 * MARGIN
    avail_h = H - top - bottom_pad
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((W - w) / 2), int(top), w, h)


def table(slide, top, headers, rows, widths, *, font=11.5, head_font=11,
          row_h=Inches(0.42)):
    n_rows, n_cols = len(rows) + 1, len(headers)
    total = W - 2 * MARGIN
    shape = slide.shapes.add_table(
        n_rows, n_cols, MARGIN, int(top), total, row_h * n_rows
    )
    tbl = shape.table
    tbl.first_row = True
    for i, frac in enumerate(widths):
        tbl.columns[i].width = Emu(int(total * frac))

    for c, head in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = BAND
        cell.margin_left = cell.margin_right = Inches(0.09)
        cell.margin_top = cell.margin_bottom = Inches(0.045)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(head_font)
        r.font.bold = True
        r.font.color.rgb = TEAL
        r.font.name = "Arial"

    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = tbl.cell(r_i, c_i)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(font)
            run.font.name = "Arial"
            # last column carries the measured effect: colour it so the table
            # can be read down that column alone
            run.font.color.rgb = TEAL if c_i == len(row) - 1 else INK
            run.font.bold = c_i == 0
    return shape


# ------------------------------------------------------------------- title
s = prs.slides.add_slide(BLANK)
tf = textbox(s, MARGIN, Inches(2.05), W - 2 * MARGIN, Inches(1.5))
para(tf, "JBROWSE 2 BENCHMARKS", 13, bold=True, color=TEAL, first=True)
tf2 = textbox(s, MARGIN, Inches(2.45), W - 2 * MARGIN, Inches(1.9))
para(tf2, "What three years bought a reader", 44, bold=True, first=True)
para(tf2, "of the 2023 paper", 44, bold=True)
rule(s, Inches(4.45), width=Inches(2.4), color=TEAL, height=Pt(3))
tf3 = textbox(s, MARGIN, Inches(4.85), Inches(9.4), Inches(1.6))
para(
    tf3,
    "Current HEAD measured against v2.4.0 — the version archived and "
    "benchmarked in the 2023 Genome Biology paper — on the same corpus, "
    "the same simulation commands, and the same machine on the same day.",
    15,
    color=MUTED,
    first=True,
)
para(tf3, "", 8)
para(tf3, "Rendering, interaction, decoding, and transport.", 15, color=MUTED)

# -------------------------------------------------------------- the thesis
s = prs.slides.add_slide(BLANK)
top = slide_head(
    s,
    "It is not the GPU",
    "A record crosses six stages between file byte and pixel: decompression, "
    "parsing, an in-memory representation, a thread boundary, an encode step, a draw.",
    eyebrow="the finding behind the numbers",
)
tf = textbox(s, MARGIN, top + Inches(0.5), Inches(11.4), Inches(3.2))
para(
    tf,
    "Each of those stages was individually reasonable. None profiled as the bottleneck.",
    22,
    first=True,
    space_after=14,
)
para(
    tf,
    "The cost was not inside them but at the boundaries between them — every stage "
    "handed the next a record in a shape that stage had to rebuild.",
    22,
    space_after=14,
    bold=True,
    color=CLAY,
)
para(
    tf,
    "A profiler charges that work to the stage doing the rebuilding, not to the "
    "boundary that forced it. So no measurement of a single layer points at it, and "
    "no owner of a single layer can remove it: the shapes have to agree, so one "
    "party has to choose all of them.",
    18,
    color=MUTED,
)
footnote(s, "We maintain the decoders, the transport and the renderer, so we chose all of them.")

# --------------------------------------------------------------- technique tables
s = prs.slides.add_slide(BLANK)
top = slide_head(
    s,
    "Techniques: representation and transport",
    "Applied at every stage a record passes through, rather than at the one that profiled worst.",
    eyebrow="1 of 2 — under the renderer",
)
table(
    s,
    top,
    ["Technique", "Replaces", "Measured effect"],
    [
        [
            "Columnar typed arrays, end to end",
            "One object per feature holding all its attributes",
            "The layout the encoder copies from and the shader reads",
        ],
        [
            "Arena allocation per decoded block",
            "One small object per record or read feature",
            "A CRAM read feature falls from 64 to 19 bytes",
        ],
        [
            "Numeric decoding instead of strings\n(packed CIGAR → opcodes; genotypes interned)",
            "A string per record, or per sample per site, that a later stage re-parses",
            "Genotype-matrix prep 1.87× / 2.47× faster, output byte-identical",
        ],
        [
            "WebAssembly for innermost loops\n(libdeflate, htscodecs, clustering kernel)",
            "A JavaScript inflate and codec implementation",
            "2.5–3× on inflation alone — but only a few % of a whole query",
        ],
        [
            "Zero-copy transport across the worker",
            "A structured clone costing time proportional to payload size",
            "Shape parsed = shape moved = shape drawn; no translation step",
        ],
    ],
    [0.30, 0.31, 0.39],
    font=10.5,
    head_font=10.5,
    row_h=Inches(0.62),
)
footnote(
    s,
    "The WebAssembly row is why the distinction matters: a stage optimized alone is bounded by the stages around it.\n"
    "The strategies that compound are the ones that change what crosses a boundary — columnar layout, the arena, and the zero-copy transport those two make possible.",
)

s = prs.slides.add_slide(BLANK)
top = slide_head(
    s,
    "Techniques: what is requested, and what is drawn",
    "Two changes to the request pattern, and the residency the layout makes possible.",
    eyebrow="2 of 2 — transport and renderer",
)
table(
    s,
    top,
    ["Technique", "Replaces", "Measured effect"],
    [
        [
            "One read per byte range, shared between callers",
            "One request per caller, cancelled by the first to abandon it",
            "No refetch or spurious failure when a neighbouring block is dropped",
        ],
        [
            "Coalesce many regions into one pass",
            "One request per region",
            "25-chromosome overview: 27 range requests and 691 kB → 3 and 312 kB",
        ],
        [
            "One chunked, compressed store for many samples",
            "One BigWig per sample, each latency-bound",
            "2,504 individuals: 15,048 requests / 24.5 s → 3 requests / 0.2 s",
        ],
        [
            "Absolute genomic coordinates resident on the GPU",
            "Rendered output bound to a specific bpPerPx",
            "Pan, zoom, recolour and re-sort are shader parameter changes — no refetch",
        ],
        [
            "One shader source in Slang → WGSL + GLSL + buffer layout",
            "Two hand-maintained implementations kept in step",
            "The encoder and the shader that reads it cannot disagree",
        ],
        [
            "WebGPU compute over the resident genotype matrix",
            "A precomputed file fixing panel, metric and window when written",
            "Those three become analytical choices; CPU fallback below threshold",
        ],
    ],
    [0.30, 0.31, 0.39],
    font=10,
    head_font=10.5,
    row_h=Inches(0.54),
)
footnote(
    s,
    "Latency was the cost on the multi-sample path, not bytes: 48.4 MB over 15,048 requests against 0.22 MB over 3.",
)

# ------------------------------------------------------------------ figures
for name, title, sub, eyebrow, note in [
    (
        "cold-load.png",
        "End to end: cold load to rendered reads",
        "Navigation → render-complete, median of 6 runs after a warmup. Lower is better.",
        "result — initial load",
        "Fetch-dominated, so it understates the renderer difference. Every build parses in a worker and pulls the same bytes.",
    ),
    (
        "speedup-vs-published.png",
        "Against the version the paper benchmarked",
        "Cold-load median of v2.4.0 ÷ median of current HEAD.",
        "result — initial load",
        "Cumulative, not isolated: three years separate v2.4.0 from HEAD and almost none of it is the renderer.",
    ),
    (
        "interaction.png",
        "What an interaction costs you",
        "Time-to-content: seconds a loading indicator sits on the track before correct content is back.",
        "result — interactivity",
        "Zoom in is the current renderer's best case and pan its worst. On a pan the region is new to both, so both pay the fetch.",
    ),
    (
        "parsers.png",
        "The parser layer underneath",
        "Decode only — no browser, no GPU. 2023 release against current, both built from source at pinned tags.",
        "result — decoding",
        "BigWig is the honest exception: 1.1–1.3× at 20x and 6–22% slower above it, on 1–3 ms operations.",
    ),
]:
    s = prs.slides.add_slide(BLANK)
    top = slide_head(s, title, sub, eyebrow=eyebrow)
    picture(s, name, top + Inches(0.06))
    footnote(s, note)

# ----------------------------------------------------------------- caveats
s = prs.slides.add_slide(BLANK)
top = slide_head(
    s,
    "What these numbers are, and are not",
    eyebrow="read this before quoting a figure",
)
tf = textbox(s, MARGIN, top + Inches(0.28), Inches(11.9), Inches(4.4))
for i, (head, body) in enumerate(
    [
        (
            "The machine was not idle.",
            "The cold-load rows were taken at 1-minute load 7–28 from unrelated jobs. "
            "This repo marks anything above 4.0 unusable, so the absolute seconds are not "
            "quotable. The four builds are measured back to back within each case, so a "
            "spike lands on all of them and the ratios survive it better than the values do.",
        ),
        (
            "The zoom table has no v2.4.0 column.",
            "The loading-indicator detector matched wording v2.4.0 does not use, so it scored "
            "the 2023 build 0 ms on every zoom — a perfect result produced by the instrument "
            "missing. Fixed, then found to break the 4.3.0 column the other way. The column is "
            "owed, not delivered.",
        ),
        (
            "Cumulative, not isolated.",
            "Three years separate v2.4.0 from HEAD. The v4.3.0 column isolates the current "
            "release; the v2.4.0 column tells a reader of the paper what the whole period bought them.",
        ),
        (
            "One machine, one locus, one workload family.",
            "Alignment pileups over a 19 kb window, against the paper's 10 kb. Same corpus and "
            "the same pbsim and wgsim commands its methods describe.",
        ),
    ]
):
    para(tf, head, 17, bold=True, color=CLAY, first=(i == 0), space_after=4)
    para(tf, body, 14, color=INK, space_after=15)

prs.save(OUT)
print(f"wrote {OUT} ({OUT.stat().st_size // 1024} kB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
